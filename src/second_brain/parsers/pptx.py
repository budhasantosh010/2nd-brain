"""PowerPoint parser with slide-level locators."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from second_brain.models import ParsedSegment
from second_brain.parsers.base import BaseParser


class PPTXParser(BaseParser):
    extensions = frozenset({".pptx"})

    def parse(self, path: Path, source_id: str):  # type: ignore[no-untyped-def]
        deck = Presentation(str(path))
        segments: list[ParsedSegment] = []
        for slide_number, slide in enumerate(deck.slides, start=1):
            parts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = str(shape.text).strip()
                    if text:
                        parts.append(text)
            segments.append(
                ParsedSegment(
                    segment_id=f"{source_id}:seg:{slide_number - 1}",
                    text="\n".join(parts),
                    locator=f"slide {slide_number}",
                    position=slide_number - 1,
                )
            )
        return self.document(
            path,
            source_id,
            segments,
            mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            metadata={"slides": len(deck.slides)},
        )
