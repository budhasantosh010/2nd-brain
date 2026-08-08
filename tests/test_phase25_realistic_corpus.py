from __future__ import annotations

from email.message import EmailMessage
from pathlib import Path

from conftest import StaticProvider
from docx import Document
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation
from pypdf import PdfWriter

from second_brain.backup import create_backup, verify_backup
from second_brain.config import load_config
from second_brain.ingest.service import IngestionService
from second_brain.knowledge.compiler import KnowledgeCompiler
from second_brain.knowledge.graph_materializer import MarkdownGraphMaterializer
from second_brain.knowledge.maps import MapGenerator
from second_brain.knowledge.projects import ProjectService, ProjectSpec, ProjectStateInput
from second_brain.knowledge.restructuring import RestructuringService
from second_brain.models import ProcessingState
from second_brain.paths import BrainPaths
from second_brain.rebuild import RebuildService
from second_brain.retrieval.service import RetrievalService
from second_brain.storage.sqlite import SQLiteStore
from second_brain.verification.consistency import ConsistencyVerifier
from second_brain.verification.service import REFUSAL, VerificationService


def _text_pdf(path: Path, text: str) -> None:
    """Write a tiny standards-compliant one-page Helvetica PDF with extractable text."""
    safe = text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
    stream = f"BT /F1 12 Tf 72 720 Td ({safe}) Tj ET".encode("latin-1")
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 5 0 R >> >> /Contents 4 0 R >>",
        b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    data = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for index, obj in enumerate(objects, start=1):
        offsets.append(len(data))
        data.extend(f"{index} 0 obj\n".encode())
        data.extend(obj)
        data.extend(b"\nendobj\n")
    xref = len(data)
    data.extend(f"xref\n0 {len(objects) + 1}\n".encode())
    data.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        data.extend(f"{offset:010d} 00000 n \n".encode())
    data.extend(
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF\n".encode()
    )
    path.write_bytes(bytes(data))


def _ingest(service: IngestionService, path: Path) -> str:
    result = service.ingest_file(path)
    assert result.source_id is not None
    assert result.raw_path is not None and result.raw_path.is_file()
    assert result.state in {ProcessingState.NEEDS_AI, ProcessingState.NEEDS_ENRICHMENT}
    return result.source_id


def _compile(paths: BrainPaths, store: SQLiteStore, source_id: str, payload: dict[str, object]):
    return KnowledgeCompiler(paths, load_config(paths), store, StaticProvider(payload)).compile_source(source_id)


def _logical_snapshot(store: SQLiteStore, project_id: str) -> dict[str, object]:
    with store.connect() as conn:
        concepts = [tuple(row) for row in conn.execute("SELECT id,title,summary FROM concepts ORDER BY id").fetchall()]
        decisions = [tuple(row) for row in conn.execute("SELECT id,status,supersedes,superseded_by FROM decisions ORDER BY id").fetchall()]
        claims = [tuple(row) for row in conn.execute("SELECT id,statement,status,source_id,supersedes,superseded_by FROM claims ORDER BY id").fetchall()]
        relationships = [tuple(row) for row in conn.execute("SELECT id,from_id,to_id,relation,source_id,provisional FROM relationships ORDER BY id").fetchall()]
        state = tuple(
            conn.execute(
                "SELECT current_state,next_action FROM project_states WHERE project_id=? AND active=1",
                (project_id,),
            ).fetchone()
        )
        gaps = [tuple(row) for row in conn.execute("SELECT id,status,resolution_id FROM questions ORDER BY id").fetchall()]
        source_count = int(conn.execute("SELECT COUNT(*) FROM sources").fetchone()[0])
        conflict_count = int(conn.execute("SELECT COUNT(*) FROM conflicts").fetchone()[0])
        fts = [tuple(row) for row in conn.execute("SELECT object_id,object_type,title,text,source_id,locator FROM search_fts ORDER BY object_id").fetchall()]
        vectors = [tuple(row) for row in conn.execute("SELECT object_id,object_type,title,text,source_id,metadata_json FROM vector_items ORDER BY object_id").fetchall()]
    return {
        "concepts": concepts,
        "decisions": decisions,
        "claims": claims,
        "relationships": relationships,
        "state": state,
        "gaps": gaps,
        "source_count": source_count,
        "conflict_count": conflict_count,
        "fts": fts,
        "vectors": vectors,
    }


def test_realistic_phase25_corpus_survives_complete_lifecycle(
    isolated_brain: BrainPaths,
    store: SQLiteStore,
    input_dir: Path,
) -> None:
    # Match the real `second-brain init` lifecycle so tracked skills/maps are indexed before use.
    RebuildService(isolated_brain).rebuild()
    store.initialize()
    service = IngestionService(isolated_brain, load_config(isolated_brain), store)
    corpus = input_dir / "Realistic Corpus With Spaces"
    corpus.mkdir(parents=True)

    long_chat = corpus / "long-ai-conversation.md"
    long_chat.write_text(
        "# AI Research Conversation\n\n" + "Hybrid retrieval must preserve evidence and current state.\n" * 300,
        encoding="utf-8",
    )
    old_source = _ingest(service, long_chat)

    handoff = corpus / "project-handoff.md"
    handoff.write_text(
        "# Handoff\n\nProject Atlas is in acceptance validation. Next action: publish the hardening branch.\n",
        encoding="utf-8",
    )
    handoff_source = _ingest(service, handoff)

    pdf = corpus / "text research.pdf"
    _text_pdf(pdf, "Text PDF evidence says acquisition cost is rising.")
    pdf_source = _ingest(service, pdf)
    assert pdf_source

    scan = corpus / "scanned research.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=300, height=300)
    with scan.open("wb") as handle:
        writer.write(handle)
    scan_result = service.ingest_file(scan)
    assert scan_result.state == ProcessingState.NEEDS_ENRICHMENT

    docx = corpus / "research.docx"
    document = Document()
    document.add_heading("DOCX Research", level=1)
    document.add_paragraph("Retention remains stable in the current cohort.")
    document.save(str(docx))
    _ingest(service, docx)

    pptx = corpus / "research.pptx"
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[1])
    slide.shapes.title.text = "Launch Review"
    slide.placeholders[1].text = "The deployment was postponed for reliability work."
    deck.save(str(pptx))
    _ingest(service, pptx)

    xlsx = corpus / "metrics.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Metrics"
    sheet.append(["metric", "value"])
    sheet.append(["CAC", 145])
    workbook.save(xlsx)
    workbook.close()
    _ingest(service, xlsx)

    html = corpus / "capture.html"
    html.write_text(
        "<html><head><title>Capture</title></head><body><p>Customer churn remained flat.</p></body></html>",
        encoding="utf-8",
    )
    _ingest(service, html)

    email = corpus / "update.eml"
    message = EmailMessage()
    message["Subject"] = "Project Atlas Update"
    message["From"] = "sender@example.test"
    message["To"] = "receiver@example.test"
    message.set_content("The launch date moved to October 4.")
    email.write_bytes(message.as_bytes())
    new_source = _ingest(service, email)

    code = corpus / "retrieval.py"
    code.write_text("CURRENT_BRANCH = 'feature/global-brain-phase2-5-hardening'\n", encoding="utf-8")
    _ingest(service, code)

    image = corpus / "dashboard.png"
    Image.new("RGB", (40, 40)).save(image)
    image_result = service.ingest_file(image)
    assert image_result.state == ProcessingState.NEEDS_ENRICHMENT

    media = corpus / "meeting.mp3"
    media.write_bytes(b"synthetic meeting bytes")
    media_result = service.ingest_file(media)
    assert media_result.state == ProcessingState.NEEDS_ENRICHMENT

    old_payload: dict[str, object] = {
        "purpose": "old project evidence",
        "entities": [],
        "project_candidates": [],
        "claims": [
            {
                "statement": "Cloud AI is enabled for all sources.",
                "source_id": old_source,
                "source_locator": "lines 1-1",
                "confidence_state": "supported",
            }
        ],
        "decisions": [
            {
                "decision": "Launch on September 10.",
                "reasoning": "Original schedule.",
                "status": "active",
                "source_ids": [old_source],
            }
        ],
        "concepts": [
            {
                "title": "Evidence Retrieval",
                "summary": "Hybrid retrieval combines lexical and semantic evidence.",
                "status": "provisional",
                "verification_state": "provisional",
                "source_ids": [old_source],
            }
        ],
        "open_loops": [],
        "questions": [],
    }
    old_compiled = _compile(isolated_brain, store, old_source, old_payload)
    old_decision = old_compiled.decisions[0]

    new_payload: dict[str, object] = {
        "purpose": "new project evidence",
        "entities": [],
        "project_candidates": [],
        "claims": [
            {
                "statement": "Only explicitly approved sources may use cloud AI.",
                "source_id": new_source,
                "source_locator": "body",
                "confidence_state": "supported",
            }
        ],
        "decisions": [
            {
                "decision": "Launch on October 4.",
                "reasoning": "The September schedule was superseded.",
                "status": "active",
                "supersedes": old_decision,
                "source_ids": [new_source],
            }
        ],
        "concepts": [
            {
                "title": "Hybrid Search Evidence",
                "summary": "Semantic and lexical retrieval work together to find evidence.",
                "status": "provisional",
                "verification_state": "provisional",
                "source_ids": [new_source],
            }
        ],
        "open_loops": [],
        "questions": [],
    }
    _compile(isolated_brain, store, new_source, new_payload)

    projects = ProjectService(isolated_brain, store)
    project_id = projects.create(
        ProjectSpec(
            title="Project Atlas Acceptance",
            goal="Verify the entire Phase 2.5 lifecycle.",
            source_ids=[old_source, new_source, handoff_source],
        )
    )
    projects.update_state(
        project_id,
        ProjectStateInput(
            current_state="Acceptance corpus is fully ingested and compiled.",
            last_completed="Mixed-format ingestion and supersession completed.",
            next_action="Back up and rebuild generated state.",
            source_ids=[handoff_source, new_source],
            latest_verified_evidence=[handoff_source, new_source],
        ),
    )
    projects.create_handoff(project_id)

    verification = VerificationService(isolated_brain, store)
    unknown = verification.ask("What is the Project Atlas emergency recovery code RECOVERY-CODE-UNRECORDED?")
    assert unknown.answer == REFUSAL
    later = corpus / "later-gap-evidence.txt"
    later.write_text(
        "Project Atlas emergency recovery code RECOVERY-CODE-7319 is the approved recovery code.",
        encoding="utf-8",
    )
    later_source = _ingest(service, later)
    later_payload: dict[str, object] = {
        "purpose": "later evidence",
        "entities": [],
        "project_candidates": [],
        "claims": [
            {
                "statement": "Project Atlas emergency recovery code is RECOVERY-CODE-7319.",
                "source_id": later_source,
                "source_locator": "lines 1-1",
                "confidence_state": "supported",
            }
        ],
        "decisions": [],
        "concepts": [],
        "open_loops": [],
        "questions": [],
    }
    _compile(isolated_brain, store, later_source, later_payload)

    MarkdownGraphMaterializer(isolated_brain, store).materialize()
    maps = MapGenerator(isolated_brain, store).generate()
    assert maps["maps_updated"] >= 1
    assert maps["objects_mapped"] >= 1
    structural = RestructuringService(isolated_brain, load_config(isolated_brain), store).analyze()
    assert structural.metrics
    assert ConsistencyVerifier(isolated_brain, store).verify().ok

    retrieval = RetrievalService(isolated_brain, store=store)
    before_results = retrieval.search("current Project Atlas acceptance state", limit=10)
    before_hits = [hit.object_id for hit in before_results]
    before_debug = [
        (hit.object_id, hit.object_type, round(hit.score, 12), hit.metadata.get("channels"), hit.updated_at)
        for hit in before_results
    ]
    before = _logical_snapshot(store, project_id)
    backup = create_backup(isolated_brain, isolated_brain.brain / "backups" / "realistic-acceptance.zip")
    assert verify_backup(backup).ok

    RebuildService(isolated_brain).rebuild()
    rebuilt = SQLiteStore(isolated_brain.db)
    rebuilt.initialize()
    after = _logical_snapshot(rebuilt, project_id)
    after_results = RetrievalService(isolated_brain, store=rebuilt).search(
        "current Project Atlas acceptance state",
        limit=10,
    )
    after_hits = [hit.object_id for hit in after_results]
    after_debug = [
        (hit.object_id, hit.object_type, round(hit.score, 12), hit.metadata.get("channels"), hit.updated_at)
        for hit in after_results
    ]
    assert after["fts"] == before["fts"], (
        f"fts_before_only={list(set(before['fts']) - set(after['fts']))[:5]} "
        f"fts_after_only={list(set(after['fts']) - set(before['fts']))[:5]}"
    )
    assert after["vectors"] == before["vectors"], (
        f"vectors_before_only={list(set(before['vectors']) - set(after['vectors']))[:5]} "
        f"vectors_after_only={list(set(after['vectors']) - set(before['vectors']))[:5]}"
    )
    assert after == before
    assert after_hits == before_hits, f"before={before_debug} after={after_debug}"
    assert ConsistencyVerifier(isolated_brain, rebuilt).verify().ok
