from __future__ import annotations

from contextlib import suppress
from email.message import EmailMessage
from pathlib import Path

from conftest import StaticProvider
from docx import Document
from openpyxl import Workbook
from PIL import Image, PngImagePlugin
from pptx import Presentation
from pypdf import PdfWriter
from pypdf.generic import DecodedStreamObject, DictionaryObject, NameObject

from second_brain.backup import create_backup, verify_backup
from second_brain.config import load_config
from second_brain.ingest.service import IngestionService
from second_brain.knowledge.compiler import KnowledgeCompiler
from second_brain.knowledge.gaps import GapResolver
from second_brain.knowledge.projects import ProjectService, ProjectSpec, ProjectStateInput
from second_brain.knowledge.restructuring import RestructuringService
from second_brain.models import ProcessingState
from second_brain.paths import BrainPaths
from second_brain.rebuild import RebuildService
from second_brain.review.service import ReviewService
from second_brain.storage.durable import read_jsonl
from second_brain.storage.sqlite import SQLiteStore
from second_brain.verification.consistency import ConsistencyVerifier
from second_brain.verification.service import REFUSAL, VerificationService


def _text_pdf(path: Path, text: str) -> None:
    writer = PdfWriter()
    page = writer.add_blank_page(width=612, height=792)
    font = DictionaryObject(
        {
            NameObject("/Type"): NameObject("/Font"),
            NameObject("/Subtype"): NameObject("/Type1"),
            NameObject("/BaseFont"): NameObject("/Helvetica"),
        }
    )
    font_ref = writer._add_object(font)  # noqa: SLF001 - fixture uses pypdf low-level API
    resources = DictionaryObject(
        {NameObject("/Font"): DictionaryObject({NameObject("/F1"): font_ref})}
    )
    page[NameObject("/Resources")] = resources
    stream = DecodedStreamObject()
    escaped = text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
    stream.set_data(f"BT /F1 12 Tf 72 720 Td ({escaped}) Tj ET".encode("latin-1"))
    page[NameObject("/Contents")] = writer._add_object(stream)  # noqa: SLF001
    with path.open("wb") as handle:
        writer.write(handle)


def _payload(source_id: str, *, summary: str = "Phase 2.5 acceptance knowledge.") -> dict[str, object]:
    return {
        "purpose": "Phase 2.5 realistic acceptance",
        "concepts": [
            {
                "title": "Phase 2.5 Acceptance",
                "summary": summary,
                "status": "provisional",
                "verification_state": "provisional",
                "source_ids": [source_id],
                "project_ids": [],
                "tags": ["acceptance"],
            }
        ],
        "claims": [],
        "entities": [],
        "decisions": [],
        "project_candidates": [],
        "open_loops": [],
        "questions": [],
    }


def test_realistic_phase25_corpus_survives_full_lifecycle(
    isolated_brain: BrainPaths,
    store: SQLiteStore,
    input_dir: Path,
) -> None:
    config = load_config(isolated_brain)
    ingest = IngestionService(isolated_brain, config, store)

    markdown = input_dir / "long conversation.md"
    markdown.write_text(
        "# AI conversation\n\n" + "The acceptance brain preserves verified evidence and provenance.\n" * 120,
        encoding="utf-8",
    )
    text_pdf = input_dir / "text evidence.pdf"
    _text_pdf(text_pdf, "Text PDF evidence says acceptance marker PDF-TEXT-7711.")
    scan_pdf = input_dir / "scanned evidence.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=300, height=300)
    with scan_pdf.open("wb") as handle:
        writer.write(handle)

    docx_path = input_dir / "handoff source.docx"
    doc = Document()
    doc.add_paragraph("DOCX acceptance handoff evidence.")
    doc.save(docx_path)

    pptx_path = input_dir / "acceptance deck.pptx"
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[1])
    slide.shapes.title.text = "Acceptance Deck"
    slide.placeholders[1].text = "PPTX acceptance evidence."
    deck.save(pptx_path)

    xlsx_path = input_dir / "acceptance sheet.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.append(["metric", "value"])
    sheet.append(["acceptance", 1])
    workbook.save(xlsx_path)
    workbook.close()

    html_path = input_dir / "capture.html"
    html_path.write_text("<html><body><p>HTML acceptance evidence.</p></body></html>", encoding="utf-8")
    email_path = input_dir / "message.eml"
    email = EmailMessage()
    email["Subject"] = "Acceptance Email"
    email["From"] = "sender@example.test"
    email["To"] = "receiver@example.test"
    email.set_content("Email acceptance evidence.")
    email_path.write_bytes(email.as_bytes())
    code_path = input_dir / "acceptance.py"
    code_path.write_text("ACCEPTANCE_MARKER = 'CODE-5511'\n", encoding="utf-8")

    image_path = input_dir / "dashboard.png"
    png_info = PngImagePlugin.PngInfo()
    png_info.add_text("Description", "Acceptance dashboard revenue grew 42 percent.")
    Image.new("RGB", (64, 64)).save(image_path, pnginfo=png_info)
    audio_path = input_dir / "voice.mp3"
    audio_path.write_bytes(b"synthetic-audio-preserved-for-capability-fallback")
    video_path = input_dir / "meeting.mp4"
    video_path.write_bytes(b"synthetic-video-preserved-for-capability-fallback")

    corpus = [
        markdown,
        text_pdf,
        scan_pdf,
        docx_path,
        pptx_path,
        xlsx_path,
        html_path,
        email_path,
        code_path,
        image_path,
        audio_path,
        video_path,
    ]
    results = [ingest.ingest_file(path) for path in corpus]
    assert all(result.source_id for result in results)
    assert results[1].state == ProcessingState.NEEDS_AI
    assert results[2].state == ProcessingState.NEEDS_ENRICHMENT
    assert results[-1].state == ProcessingState.NEEDS_ENRICHMENT

    main_source = str(results[0].source_id)
    compiler = KnowledgeCompiler(isolated_brain, config, store, StaticProvider(_payload(main_source)))
    compiled = compiler.compile_source(main_source)
    assert compiled.created_concepts
    concept_id = compiled.created_concepts[0]

    duplicate_path = input_dir / "duplicate meaning.txt"
    duplicate = ingest.ingest_file(duplicate_path) if duplicate_path.exists() else None
    if duplicate is None:
        duplicate_path.write_text("The Phase 2.5 acceptance brain preserves verified evidence and provenance.", encoding="utf-8")
        duplicate = ingest.ingest_file(duplicate_path)
    assert duplicate.source_id is not None
    duplicate_result = KnowledgeCompiler(
        isolated_brain,
        config,
        store,
        StaticProvider(_payload(duplicate.source_id)),
    ).compile_source(duplicate.source_id)
    assert concept_id in duplicate_result.duplicate_concepts

    conflict_path = input_dir / "conflict.txt"
    conflict_path.write_text("Cloud AI policy evidence deliberately conflicts.", encoding="utf-8")
    conflict = ingest.ingest_file(conflict_path)
    assert conflict.source_id is not None
    conflict_payload = {
        "purpose": "contradiction fixture",
        "concepts": [],
        "claims": [
            {
                "statement": "Cloud AI is enabled for all sources.",
                "source_id": conflict.source_id,
                "source_locator": "lines 1-1",
                "confidence_state": "supported",
            },
            {
                "statement": "Only explicitly approved sources may use cloud AI.",
                "source_id": conflict.source_id,
                "source_locator": "lines 1-1",
                "confidence_state": "supported",
            },
        ],
        "entities": [],
        "decisions": [],
        "project_candidates": [],
        "open_loops": [],
        "questions": [],
    }
    KnowledgeCompiler(isolated_brain, config, store, StaticProvider(conflict_payload)).compile_source(
        conflict.source_id
    )

    projects = ProjectService(isolated_brain, store)
    project_id = projects.create(
        ProjectSpec(title="Phase 2.5 Acceptance Project", goal="Validate the full local-first lifecycle.")
    )
    projects.update_state(
        project_id,
        ProjectStateInput(
            current_state="Acceptance lifecycle is running.",
            last_completed="Mixed-format ingestion completed.",
            next_action="Verify rebuild equality.",
            latest_verified_evidence=[main_source],
            source_ids=[main_source],
        ),
    )
    handoff_operation = projects.create_handoff(project_id)
    assert handoff_operation.startswith("OP-")

    verifier = VerificationService(isolated_brain, store)
    unknown = verifier.ask("What is the exact ORION-CODE-NOT-YET-KNOWN-9911?")
    assert unknown.answer == REFUSAL
    with store.connect() as conn:
        gap_id = str(conn.execute("SELECT id FROM questions WHERE status='open' ORDER BY created_at DESC LIMIT 1").fetchone()[0])

    later = input_dir / "later evidence.txt"
    later.write_text("The exact ORION-CODE-NOT-YET-KNOWN-9911 is BLUE-ORBIT-42.", encoding="utf-8")
    later_result = ingest.ingest_file(later)
    assert later_result.source_id is not None
    GapResolver(isolated_brain, store).recheck(source_ids=[later_result.source_id])
    with store.connect() as conn:
        gap_status = str(conn.execute("SELECT status FROM questions WHERE id=?", (gap_id,)).fetchone()[0])
    assert gap_status in {"candidate_evidence", "resolved"}

    staged = RestructuringService(isolated_brain, config, store).generate_proposals(limit=3)
    assert (isolated_brain.vault / "07 Operations" / "Structural Audit.md").is_file()
    if staged:
        with suppress(ValueError):
            ReviewService(isolated_brain, store).approve(staged[0])

    with store.connect() as conn:
        before = {
            "sources": int(conn.execute("SELECT COUNT(*) FROM sources").fetchone()[0]),
            "concepts": int(conn.execute("SELECT COUNT(*) FROM concepts").fetchone()[0]),
            "conflicts": int(conn.execute("SELECT COUNT(*) FROM conflicts").fetchone()[0]),
            "projects": int(conn.execute("SELECT COUNT(*) FROM projects").fetchone()[0]),
            "gaps": int(conn.execute("SELECT COUNT(*) FROM questions").fetchone()[0]),
        }
    assert before["sources"] >= len(corpus) + 3
    assert before["concepts"] == 1
    assert before["conflicts"] >= 1

    backup = create_backup(isolated_brain, isolated_brain.brain / "backups" / "acceptance-corpus.zip")
    assert verify_backup(backup).ok
    gap_history_before = read_jsonl(isolated_brain.brain / "ledgers" / "knowledge-gaps.jsonl")

    RebuildService(isolated_brain).rebuild()
    rebuilt = SQLiteStore(isolated_brain.db)
    rebuilt.initialize()
    with rebuilt.connect() as conn:
        after = {
            "sources": int(conn.execute("SELECT COUNT(*) FROM sources").fetchone()[0]),
            "concepts": int(conn.execute("SELECT COUNT(*) FROM concepts").fetchone()[0]),
            "conflicts": int(conn.execute("SELECT COUNT(*) FROM conflicts").fetchone()[0]),
            "projects": int(conn.execute("SELECT COUNT(*) FROM projects").fetchone()[0]),
            "gaps": int(conn.execute("SELECT COUNT(*) FROM questions").fetchone()[0]),
        }
    assert after == before
    assert read_jsonl(isolated_brain.brain / "ledgers" / "knowledge-gaps.jsonl") == gap_history_before
    assert ConsistencyVerifier(isolated_brain, rebuilt).verify().ok
