from __future__ import annotations

import json
from email.message import EmailMessage
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pypdf import PdfWriter

from second_brain.config import load_config
from second_brain.ingest.fingerprint import sha256_file
from second_brain.ingest.service import IngestionService
from second_brain.models import ParsedDocument, ProcessingState
from second_brain.paths import BrainPaths
from second_brain.storage.sqlite import SQLiteStore


def _service(paths: BrainPaths, store: SQLiteStore) -> IngestionService:
    return IngestionService(paths, load_config(paths), store)


def _assert_supported_ingest(service: IngestionService, path: Path) -> ParsedDocument:
    result = service.ingest_file(path)
    assert result.state == ProcessingState.NEEDS_AI, result
    assert result.source_id is not None
    assert result.raw_path is not None and result.raw_path.exists()
    assert result.extracted_path is not None and result.extracted_path.exists()
    assert sha256_file(path) == sha256_file(result.raw_path)
    return ParsedDocument.model_validate_json(result.extracted_path.read_text(encoding="utf-8"))


def test_text_markdown_csv_json_yaml_html_and_code(
    isolated_brain: BrainPaths, store: SQLiteStore, input_dir: Path
) -> None:
    files = {
        "sample.txt": "alpha exact-marker-123\nsecond line\n",
        "sample.md": "# Heading\nMarkdown knowledge marker.\n",
        "sample.csv": "name,value\nalpha,1\nbeta,2\n",
        "sample.tsv": "name\tvalue\nalpha\t1\n",
        "sample.json": json.dumps({"project": "brain", "status": "current"}),
        "sample.yaml": "project: brain\nstatus: current\n",
        "sample.html": "<html><head><title>HTML Source</title></head><body><h1>Section</h1><p>web evidence</p></body></html>",
        "sample.py": "BRANCH = 'feature/test-exact-id'\n",
    }
    service = _service(isolated_brain, store)
    for name, content in files.items():
        path = input_dir / name
        path.write_text(content, encoding="utf-8")
        document = _assert_supported_ingest(service, path)
        assert document.segments
        assert document.source_id.startswith("SRC-")


def test_docx_pptx_xlsx_and_eml(
    isolated_brain: BrainPaths, store: SQLiteStore, input_dir: Path
) -> None:
    docx_path = input_dir / "document.docx"
    document = Document()
    document.add_heading("Document Heading", level=1)
    document.add_paragraph("DOCX source evidence marker")
    document.save(str(docx_path))

    pptx_path = input_dir / "slides.pptx"
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[1])
    slide.shapes.title.text = "Slide Evidence"
    if slide.placeholders[1].text_frame is not None:
        slide.placeholders[1].text = "PPTX source evidence marker"
    deck.save(str(pptx_path))

    xlsx_path = input_dir / "sheet.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Evidence"
    sheet.append(["name", "value"])
    sheet.append(["xlsx-marker", 42])
    workbook.save(xlsx_path)
    workbook.close()

    eml_path = input_dir / "message.eml"
    message = EmailMessage()
    message["Subject"] = "Synthetic Email"
    message["From"] = "sender@example.test"
    message["To"] = "receiver@example.test"
    message.set_content("Email body evidence marker")
    eml_path.write_bytes(message.as_bytes())

    service = _service(isolated_brain, store)
    docs = [_assert_supported_ingest(service, path) for path in (docx_path, pptx_path, xlsx_path, eml_path)]
    assert any("DOCX source evidence" in doc.text for doc in docs)
    assert any("PPTX source evidence" in doc.text for doc in docs)
    assert any("xlsx-marker" in doc.text for doc in docs)
    assert any("Email body evidence" in doc.text for doc in docs)
    assert any("slide 1" in segment.locator for segment in docs[1].segments)
    assert any("Evidence!A1" in segment.locator for segment in docs[2].segments)


def test_pdf_is_preserved_and_scanned_pdf_need_is_explicit(
    isolated_brain: BrainPaths, store: SQLiteStore, input_dir: Path
) -> None:
    pdf_path = input_dir / "blank-scan.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=300, height=300)
    with pdf_path.open("wb") as handle:
        writer.write(handle)
    result = _service(isolated_brain, store).ingest_file(pdf_path)
    assert result.state == ProcessingState.NEEDS_ENRICHMENT
    assert result.raw_path is not None and result.raw_path.exists()
    assert result.extracted_path is not None and result.extracted_path.exists()
    assert sha256_file(pdf_path) == sha256_file(result.raw_path)
    document = ParsedDocument.model_validate_json(result.extracted_path.read_text(encoding="utf-8"))
    assert document.metadata["scanned_likely"] is True
    assert document.metadata["requires_ocr"] is True
    source_record = isolated_brain.records / f"{document.source_id}.md"
    assert "OCR/vision is required but was not run" in source_record.read_text(encoding="utf-8")


def test_folder_import_is_recursive_but_ignores_build_and_dependency_dirs(
    isolated_brain: BrainPaths, store: SQLiteStore, input_dir: Path
) -> None:
    folder = input_dir / "folder import"
    (folder / "nested").mkdir(parents=True)
    (folder / "node_modules").mkdir()
    (folder / ".git").mkdir()
    (folder / "build").mkdir()
    (folder / "a.txt").write_text("folder-a", encoding="utf-8")
    (folder / "nested" / "b.md").write_text("folder-b", encoding="utf-8")
    (folder / "node_modules" / "ignored.txt").write_text("ignore", encoding="utf-8")
    (folder / ".git" / "ignored.txt").write_text("ignore", encoding="utf-8")
    (folder / "build" / "ignored.txt").write_text("ignore", encoding="utf-8")
    results = _service(isolated_brain, store).ingest(folder)
    assert len(results) == 2
    assert {result.input_path.name for result in results} == {"a.txt", "b.md"}


def test_exact_duplicate_is_idempotent_and_same_filename_new_content_is_new_source(
    isolated_brain: BrainPaths, store: SQLiteStore, input_dir: Path
) -> None:
    service = _service(isolated_brain, store)
    first_path = input_dir / "same.txt"
    first_path.write_text("version one", encoding="utf-8")
    first = service.ingest_file(first_path)
    duplicate = service.ingest_file(first_path)
    assert duplicate.state == ProcessingState.DUPLICATE
    assert duplicate.source_id == first.source_id

    other_dir = input_dir / "other"
    other_dir.mkdir()
    second_path = other_dir / "same.txt"
    second_path.write_text("version two", encoding="utf-8")
    second = service.ingest_file(second_path)
    assert second.source_id != first.source_id
    with store.connect() as conn:
        assert conn.execute("SELECT COUNT(*) FROM sources").fetchone()[0] == 2


def test_deterministic_ingestion_populates_fts_and_vector_index_without_ai(
    isolated_brain: BrainPaths, store: SQLiteStore, input_dir: Path
) -> None:
    source = input_dir / "retrieval.txt"
    source.write_text("Exact branch marker feature/vector-index-test and retrieval architecture evidence.", encoding="utf-8")
    result = _service(isolated_brain, store).ingest_file(source)
    assert result.state == ProcessingState.NEEDS_AI
    fts = store.search_fts("feature vector index test", limit=10)
    assert any(row["source_id"] == result.source_id for row in fts)
    with store.connect() as conn:
        vector_count = conn.execute(
            "SELECT COUNT(*) FROM vector_items WHERE source_id = ?", (result.source_id,)
        ).fetchone()[0]
    assert vector_count >= 1
